# -*- coding: utf-8 -*-
"""
PhotoSlideMaker v2 — 폴더 기반 사진 슬라이드 자동 생성기
=========================================================
업로드된 사진을 폴더(서브 디렉터리)별로 그룹핑하여,
표준 양식 PPT 슬라이드의 소제목에 폴더명을 자동 매핑하고
슬라이드당 6장씩 꽉 차게 배치합니다.

주요 기능:
  - ZIP / 개별 이미지 / 폴더 이미지 모두 지원
  - 서브 폴더명 → 슬라이드 소제목 자동 매핑
  - 표준 양식(13.333×7.5 인치) 기반 본문 영역 꽉 채우기
  - 이미지 비율 보존 + 자동 크기 조정
"""

import os
import re
import zipfile
import tempfile
from collections import OrderedDict
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ------------------------------------------------------------------ #
#  이미지 파일 확장자
# ------------------------------------------------------------------ #
IMAGE_EXTENSIONS = ('.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.tif', '.tiff')


def _is_image(filename):
    return filename.lower().endswith(IMAGE_EXTENSIONS) and not os.path.basename(filename).startswith('._')


# ------------------------------------------------------------------ #
#  1. 이미지 수집: 폴더별로 그룹핑
# ------------------------------------------------------------------ #
def collect_images_by_folder(src_path, temp_dir):
    """
    ZIP 또는 디렉터리에서 이미지를 수집하여 폴더별로 그룹핑합니다.

    반환: OrderedDict { 폴더명(str) : [이미지 경로(str), ...] }
      - 루트에 있는 이미지는 키가 빈 문자열 "" 로 들어갑니다.
      - 폴더명은 알파벳/숫자 순서로 정렬됩니다.
    """
    images_by_folder = OrderedDict()

    # ZIP 처리
    if os.path.isfile(src_path) and zipfile.is_zipfile(src_path):
        with zipfile.ZipFile(src_path, 'r') as z:
            z.extractall(temp_dir)
        walk_root = temp_dir
    elif os.path.isdir(src_path):
        walk_root = src_path
    elif os.path.isfile(src_path) and _is_image(src_path):
        images_by_folder[""] = [src_path]
        return images_by_folder
    else:
        return images_by_folder

    # ZIP 안에 단일 루트 폴더가 있는 경우 한 레벨 더 들어감
    entries = [e for e in os.listdir(walk_root) if not e.startswith('.')]
    if len(entries) == 1 and os.path.isdir(os.path.join(walk_root, entries[0])):
        walk_root = os.path.join(walk_root, entries[0])

    # 서브 폴더별 이미지 수집
    root_images = []
    sub_folders = {}

    for entry in sorted(os.listdir(walk_root)):
        full_path = os.path.join(walk_root, entry)
        if os.path.isdir(full_path):
            folder_imgs = []
            for root, _, files in os.walk(full_path):
                for f in sorted(files):
                    if _is_image(f):
                        folder_imgs.append(os.path.join(root, f))
            if folder_imgs:
                sub_folders[entry] = folder_imgs
        elif os.path.isfile(full_path) and _is_image(entry):
            root_images.append(full_path)

    # 루트 이미지가 있으면 맨 앞에
    if root_images:
        images_by_folder[""] = root_images

    # 서브 폴더 (이름 순서)
    for folder_name in sorted(sub_folders.keys()):
        images_by_folder[folder_name] = sub_folders[folder_name]

    return images_by_folder


def get_image_files(src_path, temp_dir):
    """
    하위호환: 기존 API와 동일한 단일 리스트 반환.
    """
    groups = collect_images_by_folder(src_path, temp_dir)
    all_images = []
    for imgs in groups.values():
        all_images.extend(imgs)
    return all_images


# ------------------------------------------------------------------ #
#  2. 이미지 크기 계산
# ------------------------------------------------------------------ #
def calculate_fit_dimensions(img_path, max_w, max_h):
    """
    이미지 원본 비율을 보존하면서 max_w × max_h 영역을 꽉 채우도록 크기를 계산합니다.
    crop-to-fill이 아닌 fit-to-box 방식입니다.
    """
    try:
        with Image.open(img_path) as img:
            img_w, img_h = img.size
    except Exception:
        return max_w, max_h

    if img_w <= 0 or img_h <= 0:
        return max_w, max_h

    img_aspect = img_w / img_h
    box_aspect = max_w / max_h

    if img_aspect > box_aspect:
        w = max_w
        h = max_w / img_aspect
    else:
        h = max_h
        w = max_h * img_aspect

    return w, h


# ------------------------------------------------------------------ #
#  3. 슬라이드 레이아웃 계산 (3×2 격자, 꽉 차게)
# ------------------------------------------------------------------ #
# 표준 양식 기준 좌표 (13.333 × 7.5 인치 슬라이드)
# 본문 영역: 상단 소제목(~1.0in) 아래 ~ 하단(~7.2in) / 좌우 여백 0.4in
BODY_TOP = Inches(1.05)       # 소제목 아래 본문 시작
BODY_BOTTOM = Inches(7.20)    # 본문 하단 한계
BODY_LEFT = Inches(0.40)      # 좌측 여백
BODY_RIGHT = Inches(12.93)    # 우측 한계 (13.333 - 0.4)

# 격자 간격
GAP_H = Inches(0.12)          # 가로 간격
GAP_V = Inches(0.12)          # 세로 간격


def _compute_grid_boxes(n_images, slide_w=None, slide_h=None):
    """
    이미지 수(1~6)에 따른 최적 격자 좌표를 계산합니다.
    반환: [(left, top, max_w, max_h), ...]
    """
    avail_w = BODY_RIGHT - BODY_LEFT
    avail_h = BODY_BOTTOM - BODY_TOP

    n = min(n_images, 6)
    if n <= 0:
        return []

    if n == 1:
        # 전체 영역 중앙
        w = avail_w * 0.92
        h = avail_h * 0.92
        left = BODY_LEFT + (avail_w - w) / 2
        top = BODY_TOP + (avail_h - h) / 2
        return [(left, top, w, h)]

    elif n == 2:
        # 1×2 가로 배치
        w = (avail_w - GAP_H) / 2
        h = avail_h * 0.88
        top = BODY_TOP + (avail_h - h) / 2
        return [
            (BODY_LEFT, top, w, h),
            (BODY_LEFT + w + GAP_H, top, w, h),
        ]

    elif n == 3:
        # 1×3 가로 배치
        w = (avail_w - GAP_H * 2) / 3
        h = avail_h * 0.85
        top = BODY_TOP + (avail_h - h) / 2
        return [(BODY_LEFT + (w + GAP_H) * i, top, w, h) for i in range(3)]

    elif n == 4:
        # 2×2
        w = (avail_w - GAP_H) / 2
        h = (avail_h - GAP_V) / 2
        boxes = []
        for r in range(2):
            for c in range(2):
                boxes.append((
                    BODY_LEFT + (w + GAP_H) * c,
                    BODY_TOP + (h + GAP_V) * r,
                    w, h
                ))
        return boxes

    elif n == 5:
        # 3+2 (상단 3, 하단 2 중앙정렬)
        w = (avail_w - GAP_H * 2) / 3
        h = (avail_h - GAP_V) / 2
        boxes = []
        for c in range(3):
            boxes.append((BODY_LEFT + (w + GAP_H) * c, BODY_TOP, w, h))
        indent = (avail_w - (w * 2 + GAP_H)) / 2
        for c in range(2):
            boxes.append((
                BODY_LEFT + indent + (w + GAP_H) * c,
                BODY_TOP + h + GAP_V,
                w, h
            ))
        return boxes

    else:  # 6
        # 3×2 격자 — 꽉 차게
        w = (avail_w - GAP_H * 2) / 3
        h = (avail_h - GAP_V) / 2
        boxes = []
        for r in range(2):
            for c in range(3):
                boxes.append((
                    BODY_LEFT + (w + GAP_H) * c,
                    BODY_TOP + (h + GAP_V) * r,
                    w, h
                ))
        return boxes


def add_images_to_slide(slide, images, slide_w=None, slide_h=None, max_per_slide=6):
    """
    슬라이드 1장에 이미지를 격자 배치합니다.
    """
    n = min(len(images), max_per_slide)
    if n == 0:
        return

    boxes = _compute_grid_boxes(n, slide_w, slide_h)

    for i in range(n):
        img_path = images[i]
        left, top, max_w, max_h = boxes[i]

        w, h = calculate_fit_dimensions(img_path, max_w, max_h)

        # 박스 중앙 정렬
        x = left + (max_w - w) / 2
        y = top + (max_h - h) / 2

        slide.shapes.add_picture(img_path, int(x), int(y), int(w), int(h))


# ------------------------------------------------------------------ #
#  4. 슬라이드 소제목 추가
# ------------------------------------------------------------------ #
def _add_slide_subtitle(slide, text, slide_w):
    """
    표준 양식과 동일한 스타일로 소제목을 슬라이드 본문 상단에 추가합니다.
    """
    title_left = Inches(0.40)
    title_top = Inches(0.78)
    title_w = slide_w - Inches(0.80)
    title_h = Inches(0.35)

    tb = slide.shapes.add_textbox(title_left, title_top, title_w, title_h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "맑은 고딕"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(51, 51, 51)  # 짙은 회색


# ------------------------------------------------------------------ #
#  5. 표준 양식 기반 사진 슬라이드 생성 (핵심 함수)
# ------------------------------------------------------------------ #
def _find_photo_slide_index(prs):
    """
    표준 양식 템플릿에서 '사진' 플레이스홀더 슬라이드(16번째)를 찾습니다.
    '사진' 또는 '6개 삽입' 텍스트가 포함된 슬라이드의 인덱스를 반환합니다.
    """
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text:
                t = shape.text.strip()
                if ('사진' in t and '삽입' in t) or ('6개' in t and '삽입' in t):
                    return i
    return None


def _get_slide_layout_for_photo(prs, photo_slide_idx):
    """사진 슬라이드의 레이아웃을 가져옵니다."""
    if photo_slide_idx is not None and photo_slide_idx < len(prs.slides):
        return prs.slides[photo_slide_idx].slide_layout
    # Fallback: blank layout
    if len(prs.slide_layouts) > 6:
        return prs.slide_layouts[6]
    return prs.slide_layouts[0]


def _copy_slide_decorations(src_slide, dst_slide, exclude_texts=None):
    """
    소스 슬라이드의 장식 요소(헤더, 구분선, 페이지번호 등)를 새 슬라이드로 복사합니다.
    exclude_texts에 포함된 텍스트의 shape은 제외합니다.
    """
    from copy import deepcopy
    from lxml import etree

    if exclude_texts is None:
        exclude_texts = set()

    for shape in src_slide.shapes:
        # 이미지나 차트는 복사하지 않음
        if shape.shape_type == 13:  # PICTURE
            continue

        # 플레이스홀더 텍스트는 제외
        if hasattr(shape, 'text') and shape.text:
            t = shape.text.strip()
            if any(ex in t for ex in exclude_texts):
                continue

        try:
            new_el = deepcopy(shape._element)
            dst_slide.shapes._spTree.append(new_el)
        except Exception:
            pass


def generate_photo_slides_from_template(
    image_groups, template_path, output_path,
    max_per_slide=6, log=print
):
    """
    표준 양식 PPT 템플릿을 기반으로 사진 슬라이드를 생성합니다.

    Args:
        image_groups: OrderedDict { 폴더명: [이미지 경로, ...] }
        template_path: 표준 양식 PPT 템플릿 경로
        output_path: 출력 PPTX 경로
        max_per_slide: 슬라이드당 최대 사진 수 (기본 6)
        log: 로그 콜백 함수
    """
    if not template_path or not os.path.exists(template_path):
        raise FileNotFoundError(f"템플릿 파일을 찾을 수 없습니다: {template_path}")

    prs = Presentation(template_path)
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # 사진 플레이스홀더 슬라이드 찾기
    photo_idx = _find_photo_slide_index(prs)
    if photo_idx is not None:
        log(f"[템플릿] 사진 슬라이드 발견 (슬라이드 {photo_idx + 1})")
        photo_layout = _get_slide_layout_for_photo(prs, photo_idx)
        template_photo_slide = prs.slides[photo_idx]
    else:
        log("[경고] 템플릿에서 사진 슬라이드를 찾지 못했습니다. 빈 슬라이드를 사용합니다.")
        photo_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        template_photo_slide = None

    # 전체 이미지 수 카운트
    total_images = sum(len(imgs) for imgs in image_groups.values())
    if total_images == 0:
        raise FileNotFoundError("사진 파일이 발견되지 않았습니다.")

    log(f"[정보] 총 {total_images}개 이미지, {len(image_groups)}개 폴더 그룹 감지")

    slide_count = 0

    # 폴더별로 슬라이드 생성
    for folder_name, images in image_groups.items():
        if not images:
            continue

        # 한 폴더의 이미지를 max_per_slide 단위로 분할
        chunks = []
        for i in range(0, len(images), max_per_slide):
            chunks.append(images[i:i + max_per_slide])

        for chunk_idx, chunk_images in enumerate(chunks):
            # 새 슬라이드 생성
            new_slide = prs.slides.add_slide(photo_layout)
            slide_count += 1

            # 템플릿 사진 슬라이드의 장식 요소 복사 (있으면)
            if template_photo_slide is not None:
                _copy_slide_decorations(
                    template_photo_slide, new_slide,
                    exclude_texts={'사진', '삽입', '6개'}
                )

            # 소제목 생성
            if folder_name:
                subtitle = folder_name
                if len(chunks) > 1:
                    subtitle = f"{folder_name} ({chunk_idx + 1}/{len(chunks)})"
            else:
                subtitle = f"교육 활동 사진"
                if len(chunks) > 1:
                    subtitle = f"교육 활동 사진 ({chunk_idx + 1}/{len(chunks)})"

            _add_slide_subtitle(new_slide, subtitle, slide_w)

            # 이미지 배치
            add_images_to_slide(new_slide, chunk_images, slide_w, slide_h, max_per_slide)

            log(f"[진행] 슬라이드 {slide_count} 생성 완료: '{subtitle}' ({len(chunk_images)}장)")

    # 원본 사진 플레이스홀더 슬라이드 삭제
    if photo_idx is not None:
        try:
            # 새로 추가된 슬라이드가 뒤에 붙으므로, 원본 인덱스는 그대로
            _delete_slide(prs, photo_idx)
            log(f"[정리] 원본 사진 플레이스홀더 슬라이드 삭제")
        except Exception as e:
            log(f"[경고] 원본 슬라이드 삭제 실패: {e}")

    # 저장
    out_dir = os.path.dirname(os.path.abspath(output_path))
    os.makedirs(out_dir, exist_ok=True)
    prs.save(output_path)

    log(f"[성공] 총 {slide_count}장의 사진 슬라이드가 생성되었습니다 -> {output_path}")
    return output_path


def _delete_slide(prs, index):
    """슬라이드를 인덱스로 삭제합니다."""
    slides = prs.slides
    sldIdLst = slides._sldIdLst
    sldId = list(sldIdLst)[index]
    rId = sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
    try:
        prs.part.drop_rel(rId)
    except Exception:
        pass
    sldIdLst.remove(sldId)


# ------------------------------------------------------------------ #
#  6. 레거시 API: generate_photo_slides (하위호환)
# ------------------------------------------------------------------ #
def generate_photo_slides(src_path, template_path, output_path, max_per_slide=6, log=print):
    """
    이미지 리소스(ZIP/폴더)를 파싱하여 슬라이드당 최적의 수로 배치한 사진 슬라이드 PPTX를 생성합니다.
    폴더 구조가 있으면 폴더명을 소제목으로 사용합니다.
    """
    log("[시작] 이미지 파일 파싱을 수행합니다...")

    with tempfile.TemporaryDirectory() as temp_dir:
        image_groups = collect_images_by_folder(src_path, temp_dir)

        total = sum(len(imgs) for imgs in image_groups.values())
        if total == 0:
            raise FileNotFoundError("사진 파일이 발견되지 않았습니다. 올바른 ZIP 파일 또는 이미지 폴더를 지정해 주세요.")

        log(f"[정보] 총 {total}개의 이미지가 로드되었습니다.")

        # 템플릿이 표준 양식이면 template-based 생성 사용
        if template_path and os.path.exists(template_path):
            return generate_photo_slides_from_template(
                image_groups, template_path, output_path,
                max_per_slide=max_per_slide, log=log
            )

        # 템플릿 없으면 기본 빈 프레젠테이션 생성
        log("[기본] 기본 빈 템플릿(16:9)을 사용하여 프레젠테이션을 생성합니다.")
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        slide_w = prs.slide_width
        slide_h = prs.slide_height
        blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

        slide_count = 0

        for folder_name, images in image_groups.items():
            if not images:
                continue

            for i in range(0, len(images), max_per_slide):
                chunk = images[i:i + max_per_slide]
                slide = prs.slides.add_slide(blank_layout)
                slide_count += 1

                # 소제목
                if folder_name:
                    chunks_total = (len(images) + max_per_slide - 1) // max_per_slide
                    chunk_idx = i // max_per_slide
                    subtitle = folder_name
                    if chunks_total > 1:
                        subtitle = f"{folder_name} ({chunk_idx + 1}/{chunks_total})"
                else:
                    subtitle = f"교육 활동 사진 스케치 (페이지 {slide_count})"

                # 상단 제목
                title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.2), slide_w - Inches(1.2), Inches(0.5))
                tf = title_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = f"■ {subtitle}"
                p.font.name = "맑은 고딕"
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = RGBColor(45, 108, 223)

                add_images_to_slide(slide, chunk, slide_w, slide_h, max_per_slide)

                log(f"[진행] 슬라이드 {slide_count} 생성 완료 (사진 {len(chunk)}장 배치)")

        out_dir = os.path.dirname(os.path.abspath(output_path))
        os.makedirs(out_dir, exist_ok=True)
        prs.save(output_path)

        log(f"[성공] 총 {slide_count}장의 교육 사진 슬라이드가 저장되었습니다 -> {output_path}")
        return output_path
