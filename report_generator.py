# -*- coding: utf-8 -*-
"""
================================================================================
 교육 만족도 결과보고서 자동 생성기  (Survey PPT Generator v2, 범용판)
================================================================================
사용법 요약
  - 입력 3가지:  (1) 설문 로우데이터 엑셀   (2) 설정 엑셀(설정.xlsx)   (3) 템플릿 PPT
  - 출력:        결과보고서 PPT 1개

  로우데이터(구글폼 형식: 타임스탬프 + 객관식 문항 + 주관식 문항)를 직접 읽어
  ▸ 문항별 응답 분포(막대 그래프)
  ▸ 문항별 평균 점수(요약 그래프)
  ▸ 주관식 의견
  을 자동 계산/삽입하고, "파란 박스(교육마다 바뀌는 항목)"는 설정.xlsx 값으로 채웁니다.

  GUI(파일 선택창) 또는 커맨드라인(CLI) 두 방식 모두 지원합니다.
================================================================================
"""

import os
import re
import sys
import argparse
from collections import Counter

import openpyxl
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.util import Pt
from pptx.dml.color import RGBColor


# ------------------------------------------------------------------ #
#  0. 공통 설정값 (리커트 척도 등)                                     #
# ------------------------------------------------------------------ #

# ------------------------------------------------------------------ #
#  척도(리커트) 자동 감지용 라이브러리                                 #
#  - 교육마다 척도가 달라도(5점/4점/7점, '그렇다'형/'만족'형 등)        #
#    응답값을 보고 자동으로 알맞은 척도를 골라 그래프·평균을 계산합니다. #
#  - 각 척도는 "높음 → 낮음" 순서로 라벨을 나열합니다.                  #
#    점수는 위에서부터 K, K-1, ... , 1 (K=라벨 수)로 매깁니다.          #
# ------------------------------------------------------------------ #
SCALE_FAMILIES = [
    # 5점 '그렇다'형 (기존 표준)
    ["매우 그렇다", "그렇다", "보통이다", "그렇지 않다", "매우 그렇지 않다"],
    # 5점 '만족'형
    ["매우 만족", "만족", "보통", "불만족", "매우 불만족"],
    # 5점 '우수'형
    ["매우 우수", "우수", "보통", "미흡", "매우 미흡"],
    # 5점 '도움'형
    ["매우 도움", "도움", "보통", "도움 안 됨", "전혀 도움 안 됨"],
    # 4점 '만족'형 (중립 없음)
    ["매우 만족", "만족", "불만족", "매우 불만족"],
    # 4점 '그렇다'형
    ["매우 그렇다", "그렇다", "그렇지 않다", "매우 그렇지 않다"],
    # 7점 '그렇다'형
    ["매우 그렇다", "상당히 그렇다", "약간 그렇다", "보통이다",
     "약간 그렇지 않다", "상당히 그렇지 않다", "매우 그렇지 않다"],
]

# '보통이다'와 '보통'처럼 흔한 표기 흔들림을 흡수하기 위한 동의어(정규화 키 기준)
_LABEL_ALIASES = {
    "보통이다": "보통", "보통": "보통",
    "그저그렇다": "보통",
}

# 하위호환용(과거 코드/설정에서 참조): 5점 '그렇다'형 기본값
LIKERT_ORDER = SCALE_FAMILIES[0]
LIKERT_SCORE = {
    "매우그렇다": 5, "그렇다": 4, "보통이다": 3, "보통": 3,
    "그렇지않다": 2, "매우그렇지않다": 1,
}


def _alias(norm_label):
    """정규화된 라벨을 동의어 대표값으로 치환."""
    return _LABEL_ALIASES.get(norm_label, norm_label)


def detect_scale(values):
    """
    한 열(문항)의 응답값들을 보고 알맞은 척도를 자동 선택.
    반환: (order, score_map) 또는 None
      - order: 카테고리 순서 라벨 리스트(높음→낮음)
      - score_map: 정규화라벨 -> 점수(K..1)
    판정: 비어있지 않은 응답의 60% 이상이 특정 척도 라벨 집합에 속하면 그 척도로 인정.
          여러 척도가 후보면 '실제 등장 라벨 종류가 가장 많은' 척도를 선택.
    """
    norm_vals = [_alias(_norm(v)) for v in values if v is not None and str(v).strip() != ""]
    if not norm_vals:
        return None
    best = None  # (matched_label_kinds, order, score_map)
    for order in SCALE_FAMILIES:
        norm_order = [_alias(_norm(lab)) for lab in order]
        label_set = set(norm_order)
        hits = sum(1 for v in norm_vals if v in label_set)
        if hits < max(1, int(len(norm_vals) * 0.6)):
            continue
        kinds = len({v for v in norm_vals if v in label_set})
        K = len(order)
        score_map = {}
        for i, nlab in enumerate(norm_order):
            score_map[nlab] = K - i
        cand = (kinds, order, score_map)
        if best is None or cand[0] > best[0]:
            best = cand
    if best is None:
        return None
    return best[1], best[2]

# 주관식(아쉬웠던 점 등)에서 "의견 없음"류 응답을 걸러내기 위한 패턴
#  * 짧고 사실상 '없음'만 있는 답변만 제거하고,
#    "없지만 ~ 좋았습니다"처럼 뒤에 실제 내용이 이어지는 긴 답변은 보존한다.
NO_OPINION_RE = re.compile(r"^(없|해당\s*없|특별히\s*없|무응답|무|n/?a|\-|\.)", re.I)
# 순수 '없음'류로 볼 핵심 토큰(공백/문장부호 제거 후 완전일치)
_NO_OPINION_CORE = {
    "없", "없음", "없습니다", "없어요", "없다", "없었음", "없었습니다",
    "해당없음", "해당사항없음", "특별히없음", "특별히없었음", "딱히없음",
    "무", "무응답", "na", "n/a", "-", ".", "x",
}


def _is_no_opinion(text):
    """짧고 사실상 아무 내용이 없는 '없음'류 답변인지 판정."""
    t = str(text or "").strip().strip('"').strip()
    if not t:
        return True
    core = re.sub(r"[\s.·,!?~]+", "", t).lower()
    if core in _NO_OPINION_CORE:
        return True
    # '없…'로 시작하면서 아주 짧은(≈한 어절) 답변만 제거. 긴 답변은 보존.
    if len(t) <= 12 and NO_OPINION_RE.match(t):
        return True
    return False

# 주관식 문항 제목을 '좋았던 점' / '아쉬웠던 점'으로 분류하기 위한 키워드
GOOD_HEADER_KEYWORDS = ("좋았던", "좋은", "장점", "우수", "인상", "만족스", "유익")
BAD_HEADER_KEYWORDS = ("아쉬", "개선", "보완", "불편", "단점", "미흡", "건의", "바라는")


def _is_good_header(text):
    t = str(text or "")
    return any(k in t for k in GOOD_HEADER_KEYWORDS)


def _is_bad_header(text):
    t = str(text or "")
    return any(k in t for k in BAD_HEADER_KEYWORDS)


def _norm(s):
    """공백 제거한 문자열(척도 매칭용)."""
    return re.sub(r"\s+", "", str(s or "")).strip()


# ------------------------------------------------------------------ #
#  1. 리소스 경로 (PyInstaller 대응)                                   #
# ------------------------------------------------------------------ #
def resource_path(relative_path):
    try:
        base_path = sys._MEIPASS  # PyInstaller 임시 폴더
    except Exception:
        base_path = os.path.abspath(".")
    return os.path.join(base_path, relative_path)


# ================================================================== #
#  2. 설정(설정.xlsx) 읽기                                            #
# ================================================================== #
def load_config(config_path):
    """
    설정 엑셀을 읽어 하나의 dict로 반환.
    시트 구성:
      - 기본정보 : 항목/값  (표지 라벨, 제목, 강사명, 수강인원 등)
      - 교육개요 : 항목/값  (과정명, 교육일정, 교육방식, 교육대상, 교육목표)
      - 문항설정 : 문항번호/구분/섹션라벨/요약라벨  (선택. 없으면 기본값 사용)
    """
    wb = openpyxl.load_workbook(config_path, data_only=True)

    def kv_sheet(name):
        d = {}
        if name not in wb.sheetnames:
            return d
        ws = wb[name]
        for r in range(1, ws.max_row + 1):
            k = ws.cell(row=r, column=1).value
            v = ws.cell(row=r, column=2).value
            if k is None:
                continue
            key = str(k).strip()
            if key in ("항목", "구분", ""):  # 헤더행 스킵
                continue
            d[key] = ("" if v is None else str(v)).strip() if not isinstance(v, (int, float)) else v
        return d

    basic = kv_sheet("기본정보")
    overview = kv_sheet("교육개요")

    # 문항설정(선택)
    question_cfg = []
    if "문항설정" in wb.sheetnames:
        ws = wb["문항설정"]
        for r in range(1, ws.max_row + 1):
            no = ws.cell(row=r, column=1).value
            if no is None:
                continue
            try:
                no_int = int(no)
            except (ValueError, TypeError):
                continue  # 헤더/안내 행 스킵
            question_cfg.append({
                "no": no_int,
                "type": (ws.cell(row=r, column=2).value or "").strip(),
                "section_label": (ws.cell(row=r, column=3).value or "").strip(),
                "summary_label": (ws.cell(row=r, column=4).value or "").strip(),
            })

    cfg = {
        "basic": basic,
        "overview": overview,
        "questions": question_cfg,
    }
    return cfg


# ================================================================== #
#  3. 로우데이터 파싱                                                  #
# ================================================================== #
def _read_rawtable(path):
    """로우데이터를 (headers, rows)로 읽어 반환. .csv / .xlsx(.xlsm) 모두 지원.
       rows 는 완전 빈 줄을 제외한 응답 행들의 리스트."""
    ext = os.path.splitext(path)[1].lower()
    if ext == ".csv":
        import csv
        # 한글 CSV는 utf-8-sig / cp949 둘 다 자주 쓰이므로 순차 시도
        data = None
        for enc in ("utf-8-sig", "utf-8", "cp949", "euc-kr"):
            try:
                with open(path, "r", encoding=enc, newline="") as f:
                    data = list(csv.reader(f))
                break
            except (UnicodeDecodeError, UnicodeError):
                continue
        if data is None:
            with open(path, "r", encoding="utf-8", errors="replace", newline="") as f:
                data = list(csv.reader(f))
        data = [r for r in data if any((c or "").strip() for c in r)]
        if not data:
            return [], []
        headers = [h for h in data[0]]
        rows = []
        for r in data[1:]:
            # 헤더 길이에 맞춰 패딩
            row = list(r) + [None] * (len(headers) - len(r))
            rows.append(row[:len(headers)])
        return headers, rows

    wb = openpyxl.load_workbook(path, data_only=True)
    ws = wb[wb.sheetnames[0]]
    headers = [ws.cell(row=1, column=c).value for c in range(1, ws.max_column + 1)]
    rows = []
    for r in range(2, ws.max_row + 1):
        row = [ws.cell(row=r, column=c).value for c in range(1, ws.max_column + 1)]
        if all(v is None or str(v).strip() == "" for v in row):
            continue
        rows.append(row)
    return headers, rows


def parse_raw_data(excel_path):
    """
    구글폼형 로우데이터를 읽어 문항 구조를 자동 인식.
    반환:
      {
        "n_responses": 응답자 수,
        "objective": [ {"header":..., "num":1, "counts":[c1..c5], "avg":4.9,
                        "is_instructor":bool}, ... ],   # 문항 순서대로
        "subjective": [ {"header":..., "answers":[...]}, ... ],
      }
    """
    headers, rows = _read_rawtable(excel_path)
    n_responses = len(rows)

    objective, subjective = [], []
    for ci, header in enumerate(headers):
        if header is None:
            continue
        h = str(header).strip()
        # 타임스탬프 열 스킵
        if ("타임" in h) or (h.lower() in ("timestamp", "time")):
            continue
        col_vals = [rows[i][ci] for i in range(n_responses)]
        # 이 열이 객관식(리커트)인지 자동 감지 (척도 종류도 함께 판별)
        scale = detect_scale(col_vals)

        if scale is not None:
            order, score_map = scale
            norm_order = [_alias(_norm(lab)) for lab in order]
            counts = [0] * len(order)
            score_sum, score_n = 0, 0
            for v in col_vals:
                nv = _alias(_norm(v))
                if nv in score_map:
                    idx = norm_order.index(nv)
                    counts[idx] += 1
                    score_sum += score_map[nv]
                    score_n += 1
            avg = round(score_sum / score_n, 1) if score_n else 0.0
            num = _leading_number(h)
            objective.append({
                "header": h, "num": num, "counts": counts, "avg": avg,
                "order": list(order), "n_points": len(order),
            })
        else:
            answers = []
            for v in col_vals:
                if v is None:
                    continue
                t = str(v).replace("\\n", "\n").strip()
                if t:
                    answers.append(t)
            num = _leading_number(h)
            subjective.append({"header": h, "num": num, "answers": answers})

    return {
        "n_responses": n_responses,
        "objective": objective,
        "subjective": subjective,
    }


def _leading_number(text):
    m = re.match(r"\s*(\d+)", str(text))
    return int(m.group(1)) if m else None


def filter_opinions(answers, drop_no_opinion=True):
    """주관식 답변 정리: '없음'류 제거(옵션), 앞뒤 공백/따옴표 정리."""
    out = []
    for a in answers:
        t = a.replace("\\n", "\n").strip().strip('"').strip()
        if not t:
            continue
        if drop_no_opinion and _is_no_opinion(t):
            continue
        # 여러 줄 답변은 한 줄로 합쳐 가독성 유지
        t = re.sub(r"\s*\n\s*", " ", t)
        out.append(t)
    return out


# ================================================================== #
#  4. PPT 텍스트/표 헬퍼                                              #
# ================================================================== #
def _clone_run_font(src_run, dst_run):
    sf, df = src_run.font, dst_run.font
    try:
        if sf.size is not None:
            df.size = sf.size
        if sf.bold is not None:
            df.bold = sf.bold
        if sf.name:
            df.name = sf.name
        if sf.color and sf.color.type is not None:
            if sf.color.type == 1:  # RGB
                df.color.rgb = sf.color.rgb
            else:                    # THEME
                df.color.theme_color = sf.color.theme_color
    except Exception:
        pass


def set_para_text(para, text):
    """단락 텍스트를 교체하되 첫 run의 서식을 유지."""
    runs = para.runs
    if not runs:
        r = para.add_run()
        r.text = text
        return
    runs[0].text = text
    for r in runs[1:]:
        r._r.getparent().remove(r._r)


def set_cell_text(cell, text):
    """표 셀 텍스트 교체(여러 줄 지원, 첫 run 서식 유지)."""
    tf = cell.text_frame
    lines = str(text).split("\n")
    first_p = tf.paragraphs[0]
    # 기준 서식용 run 확보
    ref_run = first_p.runs[0] if first_p.runs else None
    set_para_text(first_p, lines[0])
    # 남은 단락 삭제
    for extra in list(tf.paragraphs[1:]):
        extra._p.getparent().remove(extra._p)
    # 추가 줄
    for ln in lines[1:]:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = ln
        if ref_run is not None:
            _clone_run_font(ref_run, r)


def fill_bullets(shape, items, font_size=None):
    """
    텍스트박스를 불릿(단락) 목록으로 채움. 첫 단락 서식을 기준으로 복제.
    items가 비면 '- 별도 의견 없음' 한 줄.
    """
    from pptx.enum.text import MSO_AUTO_SIZE
    tf = shape.text_frame
    # 인원수와 무관하게 박스 안에 맞도록 "넘치면 자동 축소" 설정
    try:
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass
    ref_p = tf.paragraphs[0]
    ref_run = ref_p.runs[0] if ref_p.runs else None
    if not items:
        items = ["별도 의견 없음"]
    # 첫 단락
    set_para_text(ref_p, items[0])
    if font_size and ref_p.runs:
        ref_p.runs[0].font.size = Pt(font_size)
    for extra in list(tf.paragraphs[1:]):
        extra._p.getparent().remove(extra._p)
    for it in items[1:]:
        p = tf.add_paragraph()
        p.level = ref_p.level
        r = p.add_run()
        r.text = it
        if ref_run is not None:
            _clone_run_font(ref_run, r)
        if font_size:
            r.font.size = Pt(font_size)


def find_shapes(slide, predicate):
    return [s for s in slide.shapes if predicate(s)]


def delete_shape(shape):
    shape._element.getparent().remove(shape._element)


def delete_slide(prs, index):
    """발표 순서상 index 슬라이드를 제거(목록+관계)."""
    slides = prs.slides
    sldIdLst = slides._sldIdLst
    sldId = list(sldIdLst)[index]
    rId = sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
    try:
        prs.part.drop_rel(rId)
    except Exception:
        pass
    sldIdLst.remove(sldId)


# ================================================================== #
#  5. 차트 갱신                                                       #
# ================================================================== #
def get_chart_slides(prs):
    """차트가 있는 슬라이드를 발표 순서대로 반환 [(slide, chart), ...]."""
    out = []
    for slide in prs.slides:
        for shp in slide.shapes:
            if shp.has_chart:
                out.append((slide, shp.chart))
    return out


def set_chart(chart, categories, values, series_name="계열 1"):
    cd = CategoryChartData()
    cd.categories = list(categories)
    cd.add_series(series_name, tuple(float(v) for v in values))
    chart.replace_data(cd)


# ------------------------------------------------------------------ #
#  분포차트 축 라벨(텍스트박스) 자동 갱신                              #
#  - 템플릿의 분포차트는 카테고리축이 숨겨져 있고, 척도 라벨을         #
#    별도 텍스트박스로 배치해 둔다(맨위=최고 등급, 아래=최저 등급).    #
#  - 척도가 바뀌면(만족형·7점 등) 이 라벨도 함께 바꿔줘야 그래프와    #
#    라벨이 어긋나지 않는다.                                           #
# ------------------------------------------------------------------ #
_ALL_SCALE_NORMS = set()
for _fam in SCALE_FAMILIES:
    for _lab in _fam:
        _ALL_SCALE_NORMS.add(_alias(_norm(_lab)))


def _looks_like_scale_label(text):
    return _alias(_norm(text)) in _ALL_SCALE_NORMS


def update_scale_labels(slide, order):
    """분포차트 슬라이드의 척도 라벨 텍스트박스를 detected 척도(order: 높음→낮음)에 맞게 갱신.
       위→아래 = 최고 등급→최저 등급 이 되도록 채운다. (갱신한 박스 수 반환)"""
    boxes = []
    for s in slide.shapes:
        if not s.has_text_frame:
            continue
        t = s.text_frame.text.strip()
        if not t or t.startswith("■"):
            continue
        if _looks_like_scale_label(t):
            boxes.append(s)
    if not boxes:
        return 0
    boxes.sort(key=lambda s: s.top if s.top is not None else 0)  # 위 → 아래
    n = len(boxes)
    targets = [None] * n
    targets[0] = order[0]                 # 맨 위 = 최고 등급
    low_i = len(order) - 1                # 아래쪽 칸은 맨 아래부터 최저 등급으로
    for k in range(n - 1, 0, -1):
        targets[k] = order[low_i]
        low_i = max(0, low_i - 1)
    for s, lab in zip(boxes, targets):
        set_para_text(s.text_frame.paragraphs[0], lab)
    return n


# ================================================================== #
#  6. 라벨 기본값                                                     #
# ================================================================== #
DEFAULT_SECTION_LABELS = [
    "전반적 만족도", "추천도", "교육일정", "교육 난이도",
    "강사 만족도", "강사 만족도", "교육 운영",
]
DEFAULT_SUMMARY_LABELS = [
    "전반적 만족도", "추천도", "교육일정", "교육난이도",
    "{instructor}\n만족도 1-1", "{instructor}\n만족도 1-2", "교육운영",
]


def _clean_header_label(header):
    """문항 제목에서 번호·물음표·어미를 걷어내어 짧은 라벨로."""
    t = re.sub(r"^\s*\d+\s*[.)]\s*", "", str(header or "")).strip()
    t = re.sub(r"(에 대해|에 대하여|에 대한).*$", "", t).strip()
    t = re.sub(r"(습니까|입니까|하십니까|인가요|나요|은가요|는가요)\s*\??$", "", t).strip()
    t = t.rstrip("?.").strip()
    return t or str(header)


def build_labels(cfg, objective, instructor):
    """설정의 문항설정을 우선 적용. 없으면:
       - 표준(7문항)일 때만 기본 라벨 사용
       - 문항 수가 다르면 문항 제목에서 라벨을 자동 추출(엉뚱한 기본값 방지)."""
    qcfg = {q["no"]: q for q in cfg.get("questions", []) if q.get("type", "").startswith("객")}
    is_standard = (len(objective) == len(DEFAULT_SECTION_LABELS))
    # ★ 설정의 문항설정(번호→라벨)은 '객관식 문항 수가 실제와 일치할 때만' 신뢰한다.
    #   (다른 구성의 교육에 기존 교육용 라벨이 번호만 보고 잘못 붙는 것을 방지)
    use_cfg = (len(qcfg) == len(objective)) and len(qcfg) > 0
    sections, summaries = [], []
    for i, obj in enumerate(objective):
        no = obj["num"] or (i + 1)
        auto = _clean_header_label(obj.get("header", f"문항 {no}"))
        if use_cfg and no in qcfg and qcfg[no]["section_label"]:
            sec = qcfg[no]["section_label"]
        elif is_standard and i < len(DEFAULT_SECTION_LABELS):
            sec = DEFAULT_SECTION_LABELS[i]
        else:
            sec = auto
        if use_cfg and no in qcfg and qcfg[no]["summary_label"]:
            summ = qcfg[no]["summary_label"].replace("\\n", "\n")
        elif is_standard and i < len(DEFAULT_SUMMARY_LABELS):
            summ = DEFAULT_SUMMARY_LABELS[i]
        else:
            summ = auto
        summ = summ.replace("{instructor}", instructor or "강사")
        summ = summ.replace("{강사}", instructor or "강사")
        sections.append(sec)
        summaries.append(summ)
    return sections, summaries


# ================================================================== #
#  7. 메인 생성 함수                                                  #
# ================================================================== #
def generate_report(raw_path, config, template_path, output_path,
                    drop_no_opinion=True, log=print):
    # config 는 설정.xlsx 경로(str) 또는 값 딕셔너리(dict) 둘 다 허용한다.
    for p, name in [(raw_path, "로우데이터"), (template_path, "템플릿")]:
        if not os.path.exists(p):
            raise FileNotFoundError(f"{name} 파일을 찾을 수 없습니다: {p}")

    if isinstance(config, dict):
        cfg = config
    else:
        if not config or not os.path.exists(config):
            raise FileNotFoundError(f"설정 파일을 찾을 수 없습니다: {config}")
        cfg = load_config(config)
    # 누락 키 보정(직접 입력 등으로 일부만 채워졌을 때 대비)
    cfg.setdefault("basic", {})
    cfg.setdefault("overview", {})
    cfg.setdefault("questions", [])
    data = parse_raw_data(raw_path)
    basic = cfg["basic"]
    overview = cfg["overview"]

    n_resp = data["n_responses"]
    objective = data["objective"]
    subjective = data["subjective"]
    instructor = str(basic.get("강사명", "") or "").strip()

    total = basic.get("수강인원", n_resp)
    try:
        total = int(total)
    except Exception:
        total = n_resp
    if total <= 0:
        total = n_resp
    rate = round(n_resp / total * 100) if total else 100

    log(f"[정보] 응답 {n_resp}명 / 수강 {total}명 (응답률 {rate}%)")
    log(f"[정보] 객관식 {len(objective)}문항, 주관식 {len(subjective)}문항 인식")

    sections, summaries = build_labels(cfg, objective, instructor)

    prs = Presentation(template_path)

    # ---- (A) 차트 갱신 ----
    chart_slides = get_chart_slides(prs)
    if not chart_slides:
        log("[경고] 템플릿에서 차트를 찾지 못했습니다.")
    else:
        # 첫 차트 = 요약(평균), 나머지 = 문항별 분포
        summary_chart = chart_slides[0][1]
        set_chart(summary_chart, summaries[:len(objective)],
                  [o["avg"] for o in objective], series_name="결과")
        # 평균 라벨은 항상 소수 1자리로 표시(예: 5.0)
        try:
            dl = summary_chart.plots[0].data_labels
            dl.number_format = "0.0"
            dl.number_format_is_linked = False
        except Exception:
            pass
        # 요약차트 Y축 최댓값을 척도(5점→6, 7점→8 등)에 맞춰 자동 설정(막대 잘림 방지)
        try:
            maxK = max((o.get("n_points") or 5) for o in objective)
            va = summary_chart.value_axis
            va.minimum_scale = 0.0
            va.maximum_scale = float(maxK + 1)
        except Exception:
            pass
        log(f"[완료] 요약 차트 갱신 (평균: {[o['avg'] for o in objective]})")

        dist_charts = chart_slides[1:]
        n_obj = len(objective)
        n_slots = len(dist_charts)
        leftover_slides = []   # 문항 수가 템플릿보다 적을 때 지울 여분 슬라이드
        for i, (slide, chart) in enumerate(dist_charts):
            if i < n_obj:
                order = objective[i].get("order") or LIKERT_ORDER
                counts = objective[i]["counts"]
                # ★ 이 템플릿의 분포차트는 카테고리축이 숨겨져 있고 막대는 아래→위로
                #   그려지므로, '최고 등급'이 맨 위(=상단 라벨)에 오도록 데이터를 뒤집어 넣는다.
                #   (뒤집지 않으면 8명(매우 그렇다)이 최하단 라벨 옆에 그려지는 상하 반전 발생)
                set_chart(chart, list(reversed(order)), list(reversed(counts)))
                # 축 라벨 텍스트박스도 척도에 맞게 갱신(위=최고 등급)
                update_scale_labels(slide, order)
                log(f"[완료] 문항{i+1} 분포 차트 갱신 {counts} "
                    f"({objective[i].get('n_points', len(order))}점 척도)")
            else:
                # 이 교육엔 없는 문항 슬라이드 → 통째로 제거 예약
                leftover_slides.append(slide)

        # (A-1) 문항 수 < 템플릿: 남는 문항 슬라이드 자동 삭제 (옛 그래프 잔존 방지)
        if leftover_slides:
            for sl in leftover_slides:
                try:
                    _idx = list(prs.slides).index(sl)
                    delete_slide(prs, _idx)
                except Exception:
                    pass
            log(f"[정리] 사용하지 않는 문항 슬라이드 {len(leftover_slides)}장 자동 삭제 "
                f"(객관식 {n_obj}문항)")

        # (A-2) 문항 수 > 템플릿: 슬라이드가 부족 → 명확히 경고(수동 슬라이드 추가 필요)
        if n_obj > n_slots:
            miss = n_obj - n_slots
            log(f"[경고] 객관식 {n_obj}문항인데 템플릿의 문항 슬라이드는 {n_slots}장뿐입니다. "
                f"뒤 {miss}개 문항({', '.join(str(objective[i]['num'] or i+1) for i in range(n_slots, n_obj))})은 "
                f"슬라이드가 없어 그래프가 반영되지 않았습니다. 템플릿에 문항 슬라이드를 추가하세요.")

    # ---- (B) 표지(슬라이드 1) ----
    s = prs.slides[0]
    for shp in s.shapes:
        if not shp.has_text_frame:
            continue
        txt = shp.text_frame.text.strip()
        if txt == "결과보고서 이름":
            set_para_text(shp.text_frame.paragraphs[0],
                          str(basic.get("표지_상단라벨", "결과보고서")))
        # 큰 제목(2줄)
        if ("결과보고서" in txt and shp.name.startswith("TextBox")) or txt.startswith("SK TNS"):
            paras = shp.text_frame.paragraphs
            t1 = str(basic.get("표지_제목1", "")).strip()
            t2 = str(basic.get("표지_제목2", "")).strip()
            if t1 and len(paras) >= 1:
                set_para_text(paras[0], t1)
            if t2 and len(paras) >= 2:
                set_para_text(paras[1], t2)

    # ---- (C) 교육개요 표(슬라이드 4) ----
    ov_map = [("과정명", "과정명"), ("교육일정", "교육일정"),
              ("교육방식", "교육방식"), ("교육 대상", "교육대상"),
              ("교육 목표", "교육목표")]
    for slide in prs.slides:
        for shp in slide.shapes:
            if shp.has_table:
                tbl = shp.table
                labels = [tbl.cell(r, 0).text.strip() for r in range(len(tbl.rows))]
                if "과정명" in labels[0] if labels else False:
                    for r in range(len(tbl.rows)):
                        lab = tbl.cell(r, 0).text.strip()
                        for tlab, ckey in ov_map:
                            if lab == tlab and overview.get(ckey):
                                set_cell_text(tbl.cell(r, 1), " " + str(overview[ckey]))
                    log("[완료] 교육개요 표 갱신")

    # ---- (D) 설문구성 표(슬라이드 6): 문항 열을 로우데이터 헤더로 ----
    all_q = objective + subjective
    all_q_sorted = sorted(all_q, key=lambda q: (q["num"] if q["num"] else 999))
    for slide in prs.slides:
        for shp in slide.shapes:
            if shp.has_table:
                tbl = shp.table
                head = tbl.cell(0, 0).text.strip() if len(tbl.rows) else ""
                head2 = tbl.cell(0, 1).text.strip() if len(tbl.rows) else ""
                if head == "항목" and "문항" in head2:
                    for r in range(1, len(tbl.rows)):
                        qi = r - 1
                        if qi < len(all_q_sorted):
                            set_cell_text(tbl.cell(r, 1), all_q_sorted[qi]["header"])
                    log("[완료] 설문구성 표 갱신")

    # ---- (E) 슬라이드 7: 수강/응답/응답률 ----
    for slide in prs.slides:
        for shp in slide.shapes:
            if shp.has_text_frame and "수강인원" in shp.text_frame.text:
                _update_count_line(shp, total, n_resp, rate)
                log("[완료] 응답률 문구 갱신")

    # ---- (F) 문항별 슬라이드 헤더(■) ----
    obj_by_num = {o["num"]: o for o in objective if o["num"]}
    header_shapes = []
    for slide in prs.slides:
        for shp in slide.shapes:
            if shp.has_text_frame:
                t = shp.text_frame.text.strip()
                if t.startswith("■") and re.search(r"\d+\.", t) and not shp.has_table:
                    header_shapes.append((slide, shp, t))

    for slide, shp, t in header_shapes:
        m = re.search(r"(\d+)\s*\.", t)
        if not m:
            continue
        qnum = int(m.group(1))
        # 객관식 헤더
        if qnum in obj_by_num:
            idx = [o["num"] for o in objective].index(qnum)
            sec = sections[idx] if idx < len(sections) else ""
            qtext = obj_by_num[qnum]["header"]
            set_para_text(shp.text_frame.paragraphs[0], f"■ {sec}  {qtext}")
        else:
            # 주관식 헤더(좋았던/아쉬웠던): 라벨 유지, 문항 텍스트만 로우데이터로
            subj = next((sq for sq in subjective if sq["num"] == qnum), None)
            if subj:
                # 기존 "■ 좋았던 점 8. ..." 에서 라벨만 보존
                mlab = re.match(r"■\s*([^\d]+?)\s*\d", t)
                lab = mlab.group(1).strip() if mlab else ""
                set_para_text(shp.text_frame.paragraphs[0], f"■ {lab} {subj['header']}")

    # ---- (G) 주관식 의견 채우기 ----
    #  좋았던/아쉬 외에 장점·개선점 등 다양한 문구도 인식
    good = None
    bad = None
    for sq in subjective:
        htxt = sq["header"]
        if _is_bad_header(htxt):   # '아쉬/개선/보완...' 을 먼저 확인
            bad = filter_opinions(sq["answers"], drop_no_opinion=drop_no_opinion)
        elif _is_good_header(htxt):
            good = filter_opinions(sq["answers"], drop_no_opinion=False)
    if good is None and subjective:
        good = filter_opinions(subjective[0]["answers"], drop_no_opinion=False)
    if bad is None and len(subjective) > 1:
        bad = filter_opinions(subjective[1]["answers"], drop_no_opinion=drop_no_opinion)
    good = good or []
    bad = bad or []

    # 글자 수 많으면 자동 축소
    def auto_size(items, base=11):
        n = len(items)
        if n <= 9:
            return base
        if n <= 13:
            return 10
        return 9

    slide_w = prs.slide_width
    subj_slides = []  # (slide, good_box, bad_box)
    for slide in prs.slides:
        headers_here = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        if not (any("좋았던" in h for h in headers_here) and any("아쉬" in h for h in headers_here)):
            continue
        # 헤더(■) 도형 찾기
        good_hdr = bad_hdr = None
        for s in slide.shapes:
            if s.has_text_frame and s.text_frame.text.strip().startswith("■"):
                t = s.text_frame.text
                if "좋았던" in t:
                    good_hdr = s
                elif "아쉬" in t:
                    bad_hdr = s
        gbox = bbox = None
        if good_hdr is not None and bad_hdr is not None:
            def wide_boxes_in_band(top_min, top_max):
                res = []
                for s in slide.shapes:
                    if not s.has_text_frame:
                        continue
                    if s.text_frame.text.strip().startswith("■"):
                        continue
                    if s.width < slide_w * 0.4:          # 좁은 측면라벨/페이지번호 제외
                        continue
                    if top_min <= s.top < top_max:
                        res.append(s)
                return res
            g_cands = wide_boxes_in_band(good_hdr.top, bad_hdr.top)
            b_cands = wide_boxes_in_band(bad_hdr.top, prs.slide_height + 1)
            if g_cands:
                gbox = max(g_cands, key=lambda s: s.width * s.height)
            if b_cands:
                bbox = max(b_cands, key=lambda s: s.width * s.height)
        subj_slides.append((slide, gbox, bbox))

    if subj_slides:
        slide, gbox, bbox = subj_slides[0]
        if gbox is not None:
            fill_bullets(gbox, good, font_size=auto_size(good))
        if bbox is not None:
            fill_bullets(bbox, bad, font_size=auto_size(bad))
        log(f"[완료] 주관식 의견 갱신 (좋았던 {len(good)}건 / 아쉬웠던 {len(bad)}건)")
        # 남는 주관식 슬라이드(중복 예비 페이지)는 제거
        for extra_slide, _, _ in subj_slides[1:]:
            _idx = list(prs.slides).index(extra_slide)
            delete_slide(prs, _idx)
            log("[정리] 여분의 주관식 예비 슬라이드 삭제")

    # ---- (H) 안내용 파란 박스(주석) 삭제 ----
    removed = 0
    for slide in prs.slides:
        for shp in list(slide.shapes):
            if shp.has_text_frame:
                t = shp.text_frame.text.strip()
                if t in ("내용만 바뀜", "강사명 변경", "내용만바뀜"):
                    delete_shape(shp)
                    removed += 1
    if removed:
        log(f"[정리] 안내용 파란 박스 {removed}개 삭제")

    # ---- 저장 ----
    out_dir = os.path.dirname(os.path.abspath(output_path))
    os.makedirs(out_dir, exist_ok=True)
    prs.save(output_path)
    log(f"[성공] 결과보고서 저장 완료 → {output_path}")
    return output_path


def _update_count_line(shape, total, responses, rate):
    for p in shape.text_frame.paragraphs:
        runs = p.runs
        for i, r in enumerate(runs):
            prev = runs[i - 1].text if i > 0 else ""
            txt = r.text.strip()
            if prev.endswith("총 ") and txt.isdigit():
                r.text = str(total)
            elif prev.endswith("중 ") and txt.isdigit():
                r.text = str(responses)
            elif "%" in r.text:
                r.text = re.sub(r"\d+\s*%", f"{rate}%", r.text)



# ================================================================== #
#  8. 드래그 앤 드롭 GUI                                              #
# ================================================================== #
#  - 파이썬을 몰라도 파일을 창에 "끌어다 놓으면" 자동으로 종류를 인식합니다.
#  - tkinterdnd2 가 없으면(설치 실패 등) 자동으로 "클릭해서 파일 선택"으로 대체됩니다.
# ------------------------------------------------------------------ #
import threading
import subprocess

try:
    from tkinterdnd2 import TkinterDnD, DND_FILES
    HAS_DND = True
except Exception:
    HAS_DND = False


# ------- 색상/폰트 (플랫한 밝은 테마) ------- #
CLR_BG       = "#F3F5F9"   # 전체 배경
CLR_CARD     = "#FFFFFF"   # 카드 배경
CLR_PRIMARY  = "#2D6CDF"   # 파랑 강조
CLR_PRIMARY2 = "#1B4FB0"   # 진한 파랑(hover)
CLR_OK       = "#22A45D"   # 초록(완료)
CLR_OKBG     = "#E7F6EE"   # 연초록
CLR_DROP     = "#EAF1FF"   # 드롭존 배경
CLR_DROP_HI  = "#D5E4FF"   # 드롭존 하이라이트
CLR_MUTED    = "#8A94A6"   # 회색 안내
CLR_TEXT     = "#1F2430"   # 본문
CLR_LINE     = "#DFE4EC"   # 구분선
FONT_FACE    = "맑은 고딕"  # Windows 기본 한글 폰트

# 직접 입력(엑셀 없이) 모드에서 창에 뜨는 입력 항목: (설정키, 화면라벨)
FORM_FIELDS_BASIC = [
    ("표지_상단라벨", "표지 상단 라벨"),
    ("표지_제목1",   "표지 제목 1줄"),
    ("표지_제목2",   "표지 제목 2줄"),
    ("강사명",       "강사명"),
    ("수강인원",     "수강인원(명)"),
]
FORM_FIELDS_OVERVIEW = [
    ("과정명",     "과정명"),
    ("교육일정",   "교육일정"),
    ("교육방식",   "교육방식"),
    ("교육대상",   "교육대상"),
    ("교육목표",   "교육목표"),
]


def app_dir():
    """실행 파일(또는 스크립트)이 있는 폴더."""
    if getattr(sys, "frozen", False):
        return os.path.dirname(sys.executable)
    return os.path.dirname(os.path.abspath(__file__))


def _looks_like_config(path):
    """엑셀 내부 시트로 '설정.xlsx' 여부를 판별(파일명이 달라도 인식)."""
    try:
        wb = openpyxl.load_workbook(path, read_only=True)
        sheets = set(wb.sheetnames)
        wb.close()
        return bool({"기본정보", "교육개요", "문항설정"} & sheets)
    except Exception:
        return False


def classify_file(path):
    """드롭된 파일을 raw(로우데이터)/config(설정)/template(템플릿)으로 분류."""
    ext = os.path.splitext(path)[1].lower()
    name = os.path.basename(path).lower()
    if ext in (".pptx", ".ppt", ".potx"):
        return "template"
    if ext in (".xlsx", ".xls", ".xlsm", ".csv"):
        if ("설정" in name) or ("config" in name) or ("setting" in name):
            return "config"
        if _looks_like_config(path):
            return "config"
        return "raw"
    return None


def _parse_dnd_paths(data):
    """tkinterdnd2 드롭 문자열을 파일 경로 리스트로 파싱.
    공백이 포함된 경로는 {중괄호}로 감싸여 들어옵니다."""
    paths, token, in_brace = [], "", False
    for ch in data:
        if ch == "{":
            in_brace, token = True, ""
        elif ch == "}":
            in_brace = False
            paths.append(token)
            token = ""
        elif ch == " " and not in_brace:
            if token:
                paths.append(token)
                token = ""
        else:
            token += ch
    if token:
        paths.append(token)
    return [p.strip() for p in paths if p.strip()]


class DropSlot:
    """파일 한 개를 받는 슬롯(카드). 드롭 + 클릭 선택 모두 지원."""

    def __init__(self, parent, index, title, hint, filetypes, on_change):
        import tkinter as tk
        self.tk = tk
        self.on_change = on_change
        self.filetypes = filetypes
        self.path = None

        self.card = tk.Frame(parent, bg=CLR_CARD, highlightbackground=CLR_LINE,
                             highlightthickness=1, cursor="hand2")
        inner = tk.Frame(self.card, bg=CLR_CARD)
        inner.pack(fill="both", expand=True, padx=12, pady=10)

        top = tk.Frame(inner, bg=CLR_CARD)
        top.pack(fill="x")
        self.badge = tk.Label(top, text=str(index), width=3, font=(FONT_FACE, 11, "bold"),
                              bg=CLR_DROP, fg=CLR_PRIMARY)
        self.badge.pack(side="left")
        tk.Label(top, text="  " + title, font=(FONT_FACE, 11, "bold"),
                 bg=CLR_CARD, fg=CLR_TEXT).pack(side="left")
        self.state = tk.Label(top, text="비어 있음", font=(FONT_FACE, 9),
                              bg=CLR_CARD, fg=CLR_MUTED)
        self.state.pack(side="right")

        self.fname = tk.Label(inner, text=hint, font=(FONT_FACE, 9), anchor="w",
                              bg=CLR_CARD, fg=CLR_MUTED, wraplength=470, justify="left")
        self.fname.pack(fill="x", pady=(6, 0))

        # 클릭 → 파일 선택
        for w in (self.card, inner, top, self.fname):
            w.bind("<Button-1>", lambda e: self.browse())

        # 드래그 앤 드롭 등록
        if HAS_DND:
            for w in (self.card, inner, self.fname):
                w.drop_target_register(DND_FILES)
                w.dnd_bind("<<Drop>>", self._on_drop)
                w.dnd_bind("<<DragEnter>>", lambda e: self._hi(True))
                w.dnd_bind("<<DragLeave>>", lambda e: self._hi(False))

    def grid(self, **kw):
        self.card.grid(**kw)

    def _hi(self, on):
        self.card.configure(highlightbackground=CLR_PRIMARY if on else
                            (CLR_OK if self.path else CLR_LINE),
                            highlightthickness=2 if on else 1)

    def _on_drop(self, event):
        self._hi(False)
        paths = _parse_dnd_paths(event.data)
        if paths:
            self.on_change("drop", paths)   # 상위에서 자동 분류
        return event.action

    def browse(self):
        from tkinter import filedialog
        p = filedialog.askopenfilename(filetypes=self.filetypes)
        if p:
            self.on_change("pick_one", [p], self)

    def set_path(self, path):
        self.path = path
        if path:
            self.state.configure(text="선택됨 ✓", fg=CLR_OK)
            self.badge.configure(bg=CLR_OKBG, fg=CLR_OK)
            self.fname.configure(text=os.path.basename(path), fg=CLR_TEXT)
            self.card.configure(highlightbackground=CLR_OK, highlightthickness=1)
        else:
            self.state.configure(text="비어 있음", fg=CLR_MUTED)
            self.badge.configure(bg=CLR_DROP, fg=CLR_PRIMARY)
            self.card.configure(highlightbackground=CLR_LINE, highlightthickness=1)


class App:
    def __init__(self, root):
        import tkinter as tk
        from tkinter import ttk
        self.tk = tk
        self.ttk = ttk
        self.root = root
        self.slots = {}
        self.running = False

        root.title("교육 만족도 결과보고서 생성기 v3")
        root.configure(bg=CLR_BG)
        try:
            root.tk.call("tk", "scaling", 1.15)
        except Exception:
            pass

        # 번들된 기본 템플릿/설정 예시 (파일을 안 넣어도 대체 사용)
        self.default_template = resource_path(os.path.join("assets", "템플릿_결보표양.pptx"))
        self.default_config = resource_path(os.path.join("assets", "설정_예시.xlsx"))
        if not os.path.exists(self.default_template):
            self.default_template = None
        if not os.path.exists(self.default_config):
            self.default_config = None

        self._build()
        self._prefill_from_appdir()

    # ---------------- UI 구성 ---------------- #
    def _build(self):
        tk = self.tk
        outer = tk.Frame(self.root, bg=CLR_BG)
        outer.pack(fill="both", expand=True, padx=20, pady=18)

        # 헤더
        tk.Label(outer, text="교육 만족도 결과보고서 생성기 v3",
                 font=(FONT_FACE, 17, "bold"), bg=CLR_BG, fg=CLR_TEXT).pack(anchor="w")
        sub = ("파일 3개(① 로우데이터 · ② 설정 · ③ 템플릿)를 아래에 끌어다 놓으세요. "
               "종류는 자동으로 인식됩니다.") if HAS_DND else \
              ("아래 칸을 클릭해 파일 3개(① 로우데이터 · ② 설정 · ③ 템플릿)를 선택하세요.")
        tk.Label(outer, text=sub, font=(FONT_FACE, 10), bg=CLR_BG,
                 fg=CLR_MUTED).pack(anchor="w", pady=(2, 12))

        # 통합 드롭존 (한꺼번에 놓기)
        self.dropzone = tk.Frame(outer, bg=CLR_DROP, highlightbackground=CLR_PRIMARY,
                                 highlightthickness=2, height=86, cursor="hand2")
        self.dropzone.pack(fill="x")
        self.dropzone.pack_propagate(False)
        dz_txt = ("⬇  여기에 파일 3개를 한꺼번에 끌어다 놓으세요"
                  if HAS_DND else "📂  여기를 클릭해 파일을 한꺼번에 선택하세요")
        self.dz_label = tk.Label(self.dropzone, text=dz_txt, font=(FONT_FACE, 12, "bold"),
                                 bg=CLR_DROP, fg=CLR_PRIMARY)
        self.dz_label.pack(expand=True)
        tk.Label(self.dropzone, text="(로우데이터.xlsx · 설정.xlsx · 템플릿.pptx)",
                 font=(FONT_FACE, 9), bg=CLR_DROP, fg=CLR_MUTED).pack(pady=(0, 8))
        self.dropzone.bind("<Button-1>", lambda e: self._browse_multi())
        self.dz_label.bind("<Button-1>", lambda e: self._browse_multi())
        if HAS_DND:
            for w in (self.dropzone, self.dz_label):
                w.drop_target_register(DND_FILES)
                w.dnd_bind("<<Drop>>", self._dz_drop)
                w.dnd_bind("<<DragEnter>>", lambda e: self.dropzone.configure(bg=CLR_DROP_HI))
                w.dnd_bind("<<DragLeave>>", lambda e: self.dropzone.configure(bg=CLR_DROP))

        # 개별 슬롯 3개
        slots_wrap = tk.Frame(outer, bg=CLR_BG)
        slots_wrap.pack(fill="x", pady=(12, 0))
        slots_wrap.columnconfigure(0, weight=1)

        xls = [("Excel", "*.xlsx *.xls *.xlsm *.csv")]
        ppt = [("PowerPoint", "*.pptx *.potx")]
        self.slots["raw"] = DropSlot(slots_wrap, 1, "설문 로우데이터 (.xlsx)",
                                     "구글폼 등에서 받은 응답 원본", xls, self._slot_event)
        self.slots["config"] = DropSlot(slots_wrap, 2, "설정 파일 (설정.xlsx)",
                                        "표지·강사·수강인원 등 교육별 값", xls, self._slot_event)
        self.slots["template"] = DropSlot(slots_wrap, 3, "템플릿 (.pptx)",
                                          "결과보고서 서식 PPT (그대로 사용)", ppt, self._slot_event)
        self.slots["raw"].grid(row=0, column=0, sticky="ew", pady=(0, 8))
        self.slots["config"].grid(row=1, column=0, sticky="ew", pady=(0, 8))
        self.slots["template"].grid(row=2, column=0, sticky="ew", pady=(0, 8))

        # ── 설정값 입력 방식 (엑셀 파일 / 직접 입력) ──
        modewrap = tk.Frame(outer, bg=CLR_BG)
        modewrap.pack(fill="x", pady=(10, 0))
        tk.Label(modewrap, text="설정값 입력 방식:", font=(FONT_FACE, 10, "bold"),
                 bg=CLR_BG, fg=CLR_TEXT).pack(side="left")
        self.v_mode = tk.StringVar(value="excel")
        tk.Radiobutton(modewrap, text="엑셀 파일(설정.xlsx)", variable=self.v_mode,
                       value="excel", font=(FONT_FACE, 10), bg=CLR_BG, fg=CLR_TEXT,
                       activebackground=CLR_BG, selectcolor=CLR_CARD,
                       command=self._on_mode_change).pack(side="left", padx=(8, 0))
        tk.Radiobutton(modewrap, text="직접 입력(엑셀 안 열어도 됨)", variable=self.v_mode,
                       value="form", font=(FONT_FACE, 10), bg=CLR_BG, fg=CLR_TEXT,
                       activebackground=CLR_BG, selectcolor=CLR_CARD,
                       command=self._on_mode_change).pack(side="left", padx=(6, 0))

        # 직접 입력 폼 (기본 숨김 — '직접 입력' 선택 시 표시)
        self.form_frame = tk.Frame(outer, bg=CLR_CARD,
                                   highlightbackground=CLR_LINE, highlightthickness=1)
        self.form_vars = {}
        self._loaded_questions = []
        fpad = tk.Frame(self.form_frame, bg=CLR_CARD)
        fpad.pack(fill="x", padx=12, pady=(10, 4))
        fpad.columnconfigure(1, weight=1)
        fpad.columnconfigure(3, weight=1)
        _allfields = ([("basic", k, lab) for k, lab in FORM_FIELDS_BASIC] +
                      [("overview", k, lab) for k, lab in FORM_FIELDS_OVERVIEW])
        for i, (grp, key, lab) in enumerate(_allfields):
            r, c = divmod(i, 2)
            tk.Label(fpad, text=lab, font=(FONT_FACE, 9), bg=CLR_CARD, fg=CLR_MUTED,
                     anchor="w", width=12).grid(row=r, column=c * 2, sticky="w",
                                                pady=3, padx=(0, 4))
            var = tk.StringVar()
            self.form_vars[(grp, key)] = var
            tk.Entry(fpad, textvariable=var, font=(FONT_FACE, 10), relief="solid",
                     bd=1, highlightthickness=0).grid(row=r, column=c * 2 + 1, sticky="ew",
                                                      pady=3, padx=(0, 12), ipady=2)
        tk.Label(self.form_frame,
                 text="※ 교육목표 등 긴 문장도 그대로 입력하세요. 문항 라벨은 로우데이터에서 자동 생성됩니다.",
                 font=(FONT_FACE, 8), bg=CLR_CARD, fg=CLR_MUTED,
                 anchor="w").pack(fill="x", padx=12, pady=(0, 8))

        # 저장 위치
        outrow = tk.Frame(outer, bg=CLR_CARD, highlightbackground=CLR_LINE, highlightthickness=1)
        self.outrow = outrow
        outrow.pack(fill="x", pady=(4, 0))
        inr = tk.Frame(outrow, bg=CLR_CARD)
        inr.pack(fill="x", padx=12, pady=10)
        tk.Label(inr, text="④ 저장 위치", width=10, anchor="w", font=(FONT_FACE, 11, "bold"),
                 bg=CLR_CARD, fg=CLR_TEXT).pack(side="left")
        self.v_output = tk.StringVar()
        tk.Entry(inr, textvariable=self.v_output, font=(FONT_FACE, 10), relief="solid",
                 bd=1, highlightthickness=0).pack(side="left", fill="x", expand=True, padx=8, ipady=3)
        tk.Button(inr, text="변경", font=(FONT_FACE, 9), relief="flat", bg=CLR_DROP,
                  fg=CLR_PRIMARY, activebackground=CLR_DROP_HI, cursor="hand2",
                  command=self._pick_output).pack(side="left")

        # 옵션 + 생성 버튼
        actrow = tk.Frame(outer, bg=CLR_BG)
        actrow.pack(fill="x", pady=(12, 4))
        self.v_keep = tk.BooleanVar(value=False)
        tk.Checkbutton(actrow, text="아쉬웠던 점의 '없음'류 의견도 표시",
                       variable=self.v_keep, font=(FONT_FACE, 9), bg=CLR_BG,
                       fg=CLR_MUTED, activebackground=CLR_BG,
                       selectcolor=CLR_CARD).pack(side="left")
        self.btn = tk.Button(actrow, text="결과보고서 만들기  ▶", font=(FONT_FACE, 12, "bold"),
                             bg=CLR_PRIMARY, fg="white", activebackground=CLR_PRIMARY2,
                             activeforeground="white", relief="flat", cursor="hand2",
                             padx=22, pady=9, command=self._on_run)
        self.btn.pack(side="right")
        self.btn.bind("<Enter>", lambda e: self.btn.configure(bg=CLR_PRIMARY2))
        self.btn.bind("<Leave>", lambda e: self.btn.configure(
            bg=CLR_PRIMARY if not self.running else "#9BB4E8"))

        # 진행 로그
        logwrap = tk.Frame(outer, bg=CLR_CARD, highlightbackground=CLR_LINE, highlightthickness=1)
        logwrap.pack(fill="both", expand=True, pady=(10, 0))
        tk.Label(logwrap, text="진행 상황", font=(FONT_FACE, 9, "bold"), bg=CLR_CARD,
                 fg=CLR_MUTED, anchor="w").pack(fill="x", padx=12, pady=(8, 0))
        self.txt = tk.Text(logwrap, height=9, font=("Consolas", 9), relief="flat",
                           bg="#FBFCFE", fg="#39424F", wrap="word", padx=10, pady=8,
                           highlightthickness=0)
        self.txt.pack(fill="both", expand=True, padx=8, pady=8)
        self._append_log("파일을 준비되는 대로 끌어다 놓으면 자동으로 인식합니다.")
        if not HAS_DND:
            self._append_log("(안내) 끌어다 놓기 모듈이 없어 '클릭 선택' 방식으로 동작합니다.")

    # ---------------- 파일 처리 ---------------- #
    def _slot_event(self, kind, paths, slot=None):
        if kind == "pick_one" and slot is not None:
            # 특정 슬롯에서 직접 고른 경우: 그 슬롯에 그대로 배정
            for key, s in self.slots.items():
                if s is slot:
                    self._assign(key, paths[0], forced=True)
                    break
        else:
            self._auto_assign(paths)
        self._autofill_output()

    def _dz_drop(self, event):
        self.dropzone.configure(bg=CLR_DROP)
        self._auto_assign(_parse_dnd_paths(event.data))
        self._autofill_output()
        return event.action

    def _browse_multi(self):
        from tkinter import filedialog
        ps = filedialog.askopenfilenames(
            filetypes=[("모든 입력 파일", "*.xlsx *.xls *.xlsm *.csv *.pptx *.potx"),
                       ("Excel", "*.xlsx *.xls *.xlsm *.csv"),
                       ("PowerPoint", "*.pptx *.potx")])
        if ps:
            self._auto_assign(list(ps))
            self._autofill_output()

    def _auto_assign(self, paths):
        for p in paths:
            kind = classify_file(p)
            if kind in self.slots:
                self._assign(kind, p)
            else:
                self._append_log(f"[무시] 인식 못 함: {os.path.basename(p)}")

    def _assign(self, key, path, forced=False):
        self.slots[key].set_path(path)
        label = {"raw": "① 로우데이터", "config": "② 설정", "template": "③ 템플릿"}[key]
        self._append_log(f"[인식] {label} ← {os.path.basename(path)}")
        if key == "config":
            self._prefill_form_from_config(path)

    def _on_mode_change(self):
        if self.v_mode.get() == "form":
            self.form_frame.pack(fill="x", pady=(6, 0), before=self.outrow)
            self._append_log("[설정] '직접 입력' 모드 — 아래 칸에 값을 채우세요. (설정.xlsx 불필요)")
        else:
            self.form_frame.pack_forget()
            self._append_log("[설정] '엑셀 파일' 모드 — 설정.xlsx 값을 사용합니다.")
        self._autofill_output()

    def _prefill_form_from_config(self, path):
        """설정.xlsx 를 넣으면 그 값을 직접 입력 폼에도 미리 채운다(두 방식 연동)."""
        try:
            c = load_config(path)
        except Exception:
            return
        self._loaded_questions = c.get("questions", [])
        for (grp, key), var in self.form_vars.items():
            src = c.get("basic" if grp == "basic" else "overview", {})
            if key in src and str(src[key]).strip() and not var.get().strip():
                var.set(str(src[key]))

    def _build_config_from_form(self):
        """직접 입력 폼의 값으로 설정 딕셔너리를 만든다."""
        basic, overview = {}, {}
        for (grp, key), var in self.form_vars.items():
            (basic if grp == "basic" else overview)[key] = var.get().strip()
        return {"basic": basic, "overview": overview,
                "questions": getattr(self, "_loaded_questions", [])}

    def _autofill_output(self):
        if self.v_output.get().strip():
            return  # 사용자가 이미 지정했으면 건드리지 않음
        raw = self.slots["raw"].path
        form_mode = getattr(self, "v_mode", None) and self.v_mode.get() == "form"
        cfg = None if form_mode else self.slots["config"].path
        out_dir = os.path.dirname(raw) if raw else (
            os.path.dirname(cfg) if cfg else app_dir())
        name = "결과보고서"
        title = ""
        if form_mode:
            title = (self.form_vars[("basic", "표지_제목2")].get().strip()
                     or self.form_vars[("basic", "표지_제목1")].get().strip())
        elif cfg:
            try:
                c = load_config(cfg)["basic"]
                title = (str(c.get("표지_제목2", "")).strip()
                         or str(c.get("표지_제목1", "")).strip())
            except Exception:
                title = ""
        title = title.replace("/", " ").strip()
        if title:
            name = title if "결과보고서" in title else f"{title} 결과보고서"
        elif raw:
            name = os.path.splitext(os.path.basename(raw))[0] + "_결과보고서"
        self.v_output.set(os.path.join(out_dir, name.replace(" ", "_") + ".pptx"))

    def _pick_output(self):
        from tkinter import filedialog
        p = filedialog.asksaveasfilename(defaultextension=".pptx",
                                         filetypes=[("PowerPoint", "*.pptx")],
                                         initialfile=os.path.basename(self.v_output.get() or "결과보고서.pptx"))
        if p:
            self.v_output.set(p)

    def _prefill_from_appdir(self):
        """실행 폴더에 설정.xlsx / 템플릿*.pptx 가 있으면 자동으로 미리 채움."""
        try:
            for f in os.listdir(app_dir()):
                fp = os.path.join(app_dir(), f)
                if not os.path.isfile(fp):
                    continue
                low = f.lower()
                if low.endswith((".xlsx", ".xls")) and ("설정" in f) and not self.slots["config"].path:
                    self._assign("config", fp)
                elif low.endswith(".pptx") and (("템플릿" in f) or ("양식" in f)) and not self.slots["template"].path:
                    self._assign("template", fp)
        except Exception:
            pass

    # ---------------- 로그 (스레드 안전) ---------------- #
    def _append_log(self, msg):
        self.txt.insert("end", str(msg) + "\n")
        self.txt.see("end")

    def log(self, msg):
        # 작업 스레드에서 호출되어도 안전하게 메인 스레드로 전달
        self.root.after(0, self._append_log, str(msg))

    # ---------------- 실행 ---------------- #
    def _on_run(self):
        from tkinter import messagebox
        if self.running:
            return
        raw = self.slots["raw"].path
        tpl = self.slots["template"].path or self.default_template
        out = self.v_output.get().strip()

        if not raw:
            messagebox.showwarning("입력 필요", "① 설문 로우데이터 파일을 넣어 주세요.")
            return

        # 설정값: 방식(엑셀/직접 입력)에 따라 경로 또는 딕셔너리로 결정
        if self.v_mode.get() == "form":
            cfg = self._build_config_from_form()
            b = cfg["basic"]
            if not (b.get("표지_제목1") or b.get("표지_제목2") or b.get("강사명")):
                if not messagebox.askyesno(
                        "확인", "직접 입력 값이 거의 비어 있습니다. 그래도 진행할까요?"):
                    return
        else:
            cfg = self.slots["config"].path or self.default_config
            if not cfg:
                messagebox.showwarning("입력 필요", "② 설정 파일(설정.xlsx)을 넣어 주세요.")
                return

        if not tpl:
            messagebox.showwarning("입력 필요", "③ 템플릿 PPT를 넣어 주세요.")
            return
        if not out:
            self._autofill_output()
            out = self.v_output.get().strip()

        if self.v_mode.get() == "excel" and not self.slots["config"].path and self.default_config:
            self.log("[안내] 설정 파일을 넣지 않아 기본 예시 설정을 사용합니다. (값 확인 권장)")
        if not self.slots["template"].path and self.default_template:
            self.log("[안내] 템플릿을 넣지 않아 내장 기본 템플릿을 사용합니다.")

        self.running = True
        self.btn.configure(text="생성 중…", bg="#9BB4E8", state="disabled")
        self.txt.delete("1.0", "end")
        t = threading.Thread(target=self._worker, args=(raw, cfg, tpl, out), daemon=True)
        t.start()

    def _worker(self, raw, cfg, tpl, out):
        try:
            generate_report(raw, cfg, tpl, out,
                            drop_no_opinion=not self.v_keep.get(), log=self.log)
            self.root.after(0, self._done_ok, out)
        except Exception as e:
            self.root.after(0, self._done_err, e)

    def _done_ok(self, out):
        from tkinter import messagebox
        self.running = False
        self.btn.configure(text="결과보고서 만들기  ▶", bg=CLR_PRIMARY, state="normal")
        if messagebox.askyesno("완료", f"결과보고서가 생성되었습니다.\n\n{out}\n\n"
                                        "저장된 폴더를 열어볼까요?"):
            self._open_folder(out)

    def _done_err(self, e):
        from tkinter import messagebox
        self.running = False
        self.btn.configure(text="결과보고서 만들기  ▶", bg=CLR_PRIMARY, state="normal")
        self.log(f"[오류] {e}")
        messagebox.showerror("오류", f"생성 중 문제가 발생했습니다.\n\n{e}")

    def _open_folder(self, path):
        folder = os.path.dirname(os.path.abspath(path))
        try:
            if sys.platform.startswith("win"):
                os.startfile(folder)  # noqa
            elif sys.platform == "darwin":
                subprocess.Popen(["open", folder])
            else:
                subprocess.Popen(["xdg-open", folder])
        except Exception:
            pass


# ================================================================== #
#  9. 진입점                                                          #
# ================================================================== #
def main(argv=None):
    parser = argparse.ArgumentParser(description="교육 만족도 결과보고서 생성기 v3 v2")
    parser.add_argument("--raw", help="설문 로우데이터 엑셀")
    parser.add_argument("--config", help="설정 엑셀(설정.xlsx)")
    parser.add_argument("--template", help="템플릿 PPT")
    parser.add_argument("--output", help="출력 PPT 경로")
    parser.add_argument("--keep-no-opinion", action="store_true",
                        help="아쉬웠던 점의 '없음'류 응답도 그대로 표시")
    args = parser.parse_args(argv)

    # CLI 모드: 4개 인자가 모두 있으면 창 없이 바로 생성
    if args.raw and args.config and args.template and args.output:
        generate_report(args.raw, args.config, args.template, args.output,
                        drop_no_opinion=not args.keep_no_opinion)
        return

    # GUI 모드
    try:
        if HAS_DND:
            root = TkinterDnD.Tk()
        else:
            import tkinter as tk
            root = tk.Tk()
        root.minsize(560, 720)
        try:
            root.geometry("600x760")
        except Exception:
            pass
        App(root)
        root.mainloop()
    except Exception as e:
        print("GUI를 열 수 없습니다. CLI 인자를 사용하세요.")
        print("예:  python survey_ppt_generator_v2.py --raw 로우데이터.xlsx "
              "--config 설정.xlsx --template 템플릿.pptx --output 결과.pptx")
        print("오류:", e)


if __name__ == "__main__":
    main()
